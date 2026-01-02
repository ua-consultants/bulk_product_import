# main.py — ERP Product Import/Export Service (Render-compatible)
from fastapi import FastAPI, File, UploadFile, Form, HTTPException
from fastapi.responses import StreamingResponse
from fastapi.middleware.cors import CORSMiddleware
from openpyxl import Workbook, load_workbook
import csv
from datetime import datetime
import io
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import tempfile
import os
import uuid
import json
import shutil
from typing import List, Optional
import requests
from PIL import Image
import base64

app = FastAPI(title="ERP Product Import/Export Service")

# Enable CORS for all origins (important for browser access)
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  # Allow all origins
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# -----------------------------
# HEALTH CHECK (for wake-up)
# -----------------------------
@app.get("/")
@app.get("/health")
async def health_check():
    return {"status": "ok", "service": "ERP Import/Export", "timestamp": datetime.now().isoformat()}

# -----------------------------
# IMPORT: Excel/CSV
# -----------------------------
@app.post("/import/excel")
async def import_excel(file: UploadFile = File(...)):
    if not file.filename.endswith(('.xlsx', '.csv')):
        raise HTTPException(400, "Only .xlsx or .csv files allowed")

    contents = await file.read()
    products = []

    try:
        if file.filename.endswith('.csv'):
            text = contents.decode("utf-8-sig").splitlines()  # Handle BOM
            reader = csv.DictReader(text)

            for row in reader:
                products.append({
                    'name': str(row.get('Name', '')).strip(),
                    'category_name': str(row.get('Category', '')).strip(),
                    'subcategory_name': str(row.get('Subcategory', '')).strip(),
                    'description': str(row.get('Description', '')).strip(),
                    'price': float(row.get('Price') or 0),
                    'stock_quantity': int(row.get('Stock') or 0),
                    'sku': str(row.get('SKU', '')).strip(),
                    'cft': float(row.get('CFT') or 0),
                    'material': str(row.get('Material', '')).strip(),
                    'finish': str(row.get('Finish', '')).strip(),
                    'specifications': str(row.get('Specifications', '')).strip(),
                    'main_image': None,
                    'image_data': None
                })

        else:
            wb = load_workbook(io.BytesIO(contents), data_only=True, read_only=True)
            ws = wb.active
            headers = [str(cell.value).strip() if cell.value else '' for cell in ws[1]]

            for row in ws.iter_rows(min_row=2, values_only=True):
                if not any(row):  # Skip empty rows
                    continue
                row_data = {}
                for i, value in enumerate(row):
                    key = headers[i] if i < len(headers) else f"col_{i}"
                    row_data[key] = value

                products.append({
                    'name': str(row_data.get('Name', '')).strip(),
                    'category_name': str(row_data.get('Category', '')).strip(),
                    'subcategory_name': str(row_data.get('Subcategory', '')).strip(),
                    'description': str(row_data.get('Description', '')).strip(),
                    'price': float(row_data.get('Price') or 0),
                    'stock_quantity': int(row_data.get('Stock') or 0),
                    'sku': str(row_data.get('SKU', '')).strip(),
                    'cft': float(row_data.get('CFT') or 0),
                    'material': str(row_data.get('Material', '')).strip(),
                    'finish': str(row_data.get('Finish', '')).strip(),
                    'specifications': str(row_data.get('Specifications', '')).strip(),
                    'main_image': None,
                    'image_data': None
                })
            wb.close()

    except Exception as e:
        raise HTTPException(400, f"Failed to parse file: {str(e)}")

    return {"products": products}

# -----------------------------
# IMPORT: PowerPoint (images + text extraction)
# -----------------------------
@app.post("/import/pptx")
async def import_pptx(file: UploadFile = File(...), category_id: int = Form(0), subcategory_id: Optional[str] = Form(None)):
    if not file.filename.endswith('.pptx'):
        raise HTTPException(400, "Only .pptx files allowed")
    
    contents = await file.read()
    temp_dir = tempfile.mkdtemp()
    tmp_path = os.path.join(temp_dir, "uploaded.pptx")
    
    try:
        with open(tmp_path, "wb") as f:
            f.write(contents)

        prs = Presentation(tmp_path)
        products = []

        for slide_num, slide in enumerate(prs.slides, 1):
            # Extract all text from slide
            all_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    all_text.append(shape.text.strip())
            
            full_text = "\n".join(all_text)
            
            # Extract all images from slide
            images = []
            for shape in slide.shapes:
                if hasattr(shape, "image"):
                    img = shape.image
                    ext = img.ext or "png"
                    img_name = f"slide_{slide_num}_img_{len(images) + 1}_{uuid.uuid4().hex[:8]}.{ext}"
                    img_base64 = base64.b64encode(img.blob).decode('utf-8')
                    images.append({
                        'filename': img_name,
                        'data': img_base64
                    })
            
            # Parse text to extract product info
            product_data = parse_product_text(full_text)
            
            # Determine product name
            if product_data.get('name'):
                product_name = product_data['name']
            elif images:
                product_name = f"Product from Slide {slide_num}"
            else:
                product_name = f"Product {slide_num}"
            
            # Build product object
            product = {
                "name": product_name,
                "description": product_data.get('description', ''),
                "category_id": category_id,
                "subcategory_id": int(subcategory_id) if subcategory_id and subcategory_id.isdigit() else None,
                "price": product_data.get('price', 0.0),
                "stock_quantity": product_data.get('stock', 0),
                "sku": product_data.get('sku', ''),
                "main_image": None,
                "main_image_data": None,
                "gallery_images": [],
                "dimensions": product_data.get('dimensions', ''),
                "material": product_data.get('material', ''),
                "specifications": product_data.get('specifications', '')
            }
            
            # Handle images
            if len(images) > 0:
                # First image is main image
                product['main_image'] = images[0]['filename']
                product['main_image_data'] = images[0]['data']
                
                # Rest are gallery images
                if len(images) > 1:
                    product['gallery_images'] = images[1:]
            
            # Only add product if it has at least a name or an image
            if product['name'] or images:
                products.append(product)

        return {"products": products, "total": len(products)}

    except Exception as e:
        raise HTTPException(500, f"Failed to process PowerPoint: {str(e)}")
    finally:
        shutil.rmtree(temp_dir, ignore_errors=True)


def parse_product_text(text: str) -> dict:
    """
    Parse product information from text using keywords
    """
    import re
    
    data = {
        'name': '',
        'description': '',
        'price': 0.0,
        'stock': 0,
        'sku': '',
        'dimensions': '',
        'material': '',
        'specifications': ''
    }
    
    if not text:
        return data
    
    lines = text.split('\n')
    text_lower = text.lower()
    
    # Extract Price
    price_patterns = [
        r'(?:price|rate|cost|rs\.?|₹)\s*[:=-]?\s*(?:rs\.?|₹)?\s*([\d,]+(?:\.\d{2})?)',
        r'(?:rs\.?|₹)\s*([\d,]+(?:\.\d{2})?)',
    ]
    for pattern in price_patterns:
        match = re.search(pattern, text_lower)
        if match:
            price_str = match.group(1).replace(',', '')
            try:
                data['price'] = float(price_str)
                break
            except:
                pass
    
    # Extract SKU
    sku_patterns = [
        r'(?:sku|item\s*code|product\s*code|code)\s*[:=-]?\s*([A-Z0-9-]+)',
        r'\b([A-Z]{2,}\d{3,})\b',  # Pattern like ABC123
    ]
    for pattern in sku_patterns:
        match = re.search(pattern, text, re.IGNORECASE)
        if match:
            data['sku'] = match.group(1).strip()
            break
    
    # Extract Dimensions/Size
    dim_patterns = [
        r'(?:size|dimension|dimensions|measurements?)\s*[:=-]?\s*(.+?)(?:\n|$)',
        r'(\d+\s*[xX×]\s*\d+\s*[xX×]?\s*\d*)\s*(?:cm|mm|inch|ft|m)',
    ]
    for pattern in dim_patterns:
        match = re.search(pattern, text, re.IGNORECASE)
        if match:
            data['dimensions'] = match.group(1).strip()
            break
    
    # Extract Material
    material_patterns = [
        r'(?:material|made\s*of)\s*[:=-]?\s*(.+?)(?:\n|$)',
    ]
    for pattern in material_patterns:
        match = re.search(pattern, text, re.IGNORECASE)
        if match:
            data['material'] = match.group(1).strip()
            break
    
    # Extract Stock/Quantity
    stock_patterns = [
        r'(?:stock|quantity|qty|available)\s*[:=-]?\s*(\d+)',
    ]
    for pattern in stock_patterns:
        match = re.search(pattern, text_lower)
        if match:
            try:
                data['stock'] = int(match.group(1))
                break
            except:
                pass
    
    # Extract Description
    desc_patterns = [
        r'(?:desc|description)\s*[:=-]?\s*(.+?)(?:\n\n|$)',
    ]
    for pattern in desc_patterns:
        match = re.search(pattern, text, re.IGNORECASE | re.DOTALL)
        if match:
            data['description'] = match.group(1).strip()
            break
    
    # If no explicit description, use first few lines as description
    if not data['description'] and len(lines) > 1:
        # Skip lines that look like labels or have extracted data
        desc_lines = []
        for line in lines[:5]:  # First 5 lines
            line = line.strip()
            if len(line) > 15 and not re.match(r'^(price|sku|size|material|stock)', line, re.IGNORECASE):
                desc_lines.append(line)
        if desc_lines:
            data['description'] = ' '.join(desc_lines)[:500]  # Limit to 500 chars
    
    # Extract Product Name (first line if it's substantial)
    if lines:
        first_line = lines[0].strip()
        # First line is name if it's not a label and has reasonable length
        if len(first_line) > 3 and len(first_line) < 100:
            if not re.match(r'^(price|sku|size|material|stock|desc)', first_line, re.IGNORECASE):
                data['name'] = first_line
    
    # Compile specifications from all extracted data
    specs = []
    if data['dimensions']:
        specs.append(f"Dimensions: {data['dimensions']}")
    if data['material']:
        specs.append(f"Material: {data['material']}")
    
    if specs:
        data['specifications'] = '; '.join(specs)
    
    return data

# -----------------------------
# EXPORT: Excel
# -----------------------------
@app.post("/export/excel")
async def export_excel(products: str = Form(...)):
    try:
        data = json.loads(products)
    except Exception as e:
        raise HTTPException(400, f"Invalid JSON: {str(e)}")

    if not isinstance(data, list):
        raise HTTPException(400, "Products must be a JSON array")

    if not data:
        raise HTTPException(400, "No products to export")

    wb = Workbook(write_only=True)
    ws = wb.create_sheet("Products")

    headers = [
        'name', 'category_name', 'subcategory_name', 'description',
        'price', 'stock_quantity', 'sku', 'cft', 'material', 'finish', 'specifications'
    ]
    ws.append(headers)

    for item in data:
        row = [
            item.get('name', ''),
            item.get('category_name', ''),
            item.get('subcategory_name', ''),
            item.get('description', ''),
            float(item.get('price', 0)),
            int(item.get('stock_quantity', 0)),
            item.get('sku', ''),
            float(item.get('cft', 0)),
            item.get('material', ''),
            item.get('finish', ''),
            item.get('specifications', '')
        ]
        ws.append(row)

    output = io.BytesIO()
    wb.save(output)
    output.seek(0)

    return StreamingResponse(
        output,
        media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        headers={"Content-Disposition": "attachment; filename=products_export.xlsx"}
    )

# -----------------------------
# Helper: Download and process image
# -----------------------------
def download_image(image_url_or_path, base_url=None):
    """
    Download image from URL or base64 data
    Returns PIL Image object or None
    """
    try:
        # Check if it's base64 data
        if isinstance(image_url_or_path, str) and image_url_or_path.startswith('data:image'):
            # Extract base64 data
            header, encoded = image_url_or_path.split(',', 1)
            img_data = base64.b64decode(encoded)
            return Image.open(io.BytesIO(img_data))
        
        # Check if it's a base64 string without header
        if isinstance(image_url_or_path, str) and len(image_url_or_path) > 100 and not image_url_or_path.startswith('http'):
            try:
                img_data = base64.b64decode(image_url_or_path)
                return Image.open(io.BytesIO(img_data))
            except:
                pass
        
        # Construct full URL if base_url provided
        if base_url and not image_url_or_path.startswith('http'):
            image_url = f"{base_url.rstrip('/')}/{image_url_or_path.lstrip('/')}"
        else:
            image_url = image_url_or_path
        
        # Download from URL
        response = requests.get(image_url, timeout=10, headers={'User-Agent': 'Mozilla/5.0'})
        if response.status_code == 200:
            return Image.open(io.BytesIO(response.content))
    except Exception as e:
        print(f"Failed to download image {image_url_or_path}: {e}")
    
    return None

def resize_image_for_ppt(img, max_width=3.5, max_height=3.5):
    """
    Resize PIL image to fit within PowerPoint dimensions (in inches)
    Returns width and height in Inches
    """
    img_width, img_height = img.size
    aspect = img_width / img_height
    
    # Convert inches to pixels (assume 96 DPI)
    max_w_px = max_width * 96
    max_h_px = max_height * 96
    
    if img_width > max_w_px or img_height > max_h_px:
        if aspect > 1:  # Wider than tall
            new_width = max_w_px
            new_height = new_width / aspect
        else:  # Taller than wide
            new_height = max_h_px
            new_width = new_height * aspect
    else:
        new_width = img_width
        new_height = img_height
    
    # Convert back to inches
    return Inches(new_width / 96), Inches(new_height / 96)

# -----------------------------
# EXPORT: PowerPoint with Images
# -----------------------------
@app.post("/export/pptx")
async def export_pptx(products: str = Form(...), base_url: str = Form(None)):
    try:
        data = json.loads(products)
    except Exception as e:
        raise HTTPException(400, f"Invalid JSON: {str(e)}")

    if not isinstance(data, list):
        raise HTTPException(400, "Products must be a JSON array")

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Clear default slides
    if len(prs.slides) > 0:
        prs.slides._sldIdLst = prs.slides._sldIdLst[:0]

    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    left = top = Inches(1)
    width = Inches(11.33)
    height = Inches(2)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    tf = title_box.text_frame
    tf.text = "Product Catalog"
    tf.paragraphs[0].font.size = Pt(44)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    subtitle_box = slide.shapes.add_textbox(left, top + Inches(2.5), width, Inches(1))
    tf2 = subtitle_box.text_frame
    tf2.text = f"Generated on {datetime.now().strftime('%B %d, %Y')}"
    tf2.paragraphs[0].font.size = Pt(24)
    tf2.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Product slides: 2 per slide
    chunks = [data[i:i+2] for i in range(0, len(data), 2)]
    
    for chunk in chunks:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        for idx, prod in enumerate(chunk):
            x_base = Inches(0.5 + idx * 6.5)
            y = Inches(0.5)
            
            # Try to add product image
            image_added = False
            main_image = prod.get('main_image') or prod.get('image_url')
            
            if main_image:
                img = download_image(main_image, base_url)
                if img:
                    try:
                        # Save to temporary file
                        temp_img = io.BytesIO()
                        img_format = 'PNG' if img.mode == 'RGBA' else 'JPEG'
                        img.save(temp_img, format=img_format)
                        temp_img.seek(0)
                        
                        # Calculate size
                        width, height = resize_image_for_ppt(img, max_width=5.5, max_height=3.5)
                        
                        # Center image in its column
                        x_img = x_base + (Inches(6) - width) / 2
                        
                        # Add to slide
                        slide.shapes.add_picture(temp_img, x_img, y, width=width, height=height)
                        image_added = True
                        y += height + Inches(0.2)
                    except Exception as e:
                        print(f"Failed to add image: {e}")
            
            # If no image added, leave space
            if not image_added:
                y += Inches(2.5)
            
            # Product name
            name_box = slide.shapes.add_textbox(x_base, y, Inches(6), Inches(0.5))
            name_tf = name_box.text_frame
            name_tf.text = str(prod.get('name', 'Unnamed Product'))
            name_tf.paragraphs[0].font.bold = True
            name_tf.paragraphs[0].font.size = Pt(16)
            name_tf.word_wrap = True
            
            # Details
            y += Inches(0.6)
            details = []
            
            if prod.get('price'):
                details.append(f"Price: ${float(prod.get('price', 0)):.2f}")
            if prod.get('sku'):
                details.append(f"SKU: {prod['sku']}")
            if prod.get('stock_quantity') is not None:
                details.append(f"Stock: {prod['stock_quantity']}")
            if prod.get('cft'):
                details.append(f"CFT: {prod['cft']}")
            if prod.get('material'):
                details.append(f"Material: {prod['material']}")
            if prod.get('finish'):
                details.append(f"Finish: {prod['finish']}")
            
            details_text = '\n'.join(details)
            
            details_box = slide.shapes.add_textbox(x_base, y, Inches(6), Inches(2))
            details_tf = details_box.text_frame
            details_tf.text = details_text
            details_tf.paragraphs[0].font.size = Pt(11)
            details_tf.word_wrap = True

    output = io.BytesIO()
    prs.save(output)
    output.seek(0)

    return StreamingResponse(
        output,
        media_type='application/vnd.openxmlformats-officedocument.presentationml.presentation',
        headers={"Content-Disposition": "attachment; filename=product_catalog.pptx"}
    )
